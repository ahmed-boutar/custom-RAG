from pptx import Presentation
import fitz #PyMuPDF
import os 


class IngesterAndParser():
    def __init__(self):
        pass

    def extract_text_from_pptx(self, file_path):
        print("Extracting text from pptx file...\n")
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text.append("\n".join(slide_text))
        except Exception as e:
            print(f"Error reading PPTX file: {e}")
            return None

        print("Extraction complete.\n")
        
        res = {"filename": file_path.name, "content": text, "type": "pptx"}
        return res
    

    def process_documents(self, dir):
        documents = []
        
        for filename in os.listdir(dir):
            filename = os.path.join(dir, filename)
            print(f'Processing file: {filename}...\n')
            if filename.endswith(".pdf"):
                #TODO: Implement logic for pdf files
                pass
            elif filename.endswith(".pptx"):
                content = self.extract_text_from_pptx(filename)
                documents.append(content)
            else:
                print(f"Unsupported file format: {filename}")
        
        return documents

def main():
    DATA_DIR = "../data"
    ingesterAndParser = IngesterAndParser()
    
    extracted_text = ingesterAndParser.process_documents(DATA_DIR)
    print(extracted_text)

if __name__ == "__main__":
    main()
